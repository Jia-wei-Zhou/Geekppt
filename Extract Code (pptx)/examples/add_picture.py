from pptx import Presentation
from pptx.util import Inches
import os


img_path = "../resources" + os.path.sep + "1.jpeg"
if not os.path.exists(img_path):
    print("{} does not exist".format(img_path))
    exit(1)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(0)
height = Inches(2.5)
img = slide.shapes.add_picture(img_path, left, top, height=height)

left = top = Inches(2.5)
height = Inches(4.5)
img = slide.shapes.add_picture(img_path, left, top, height=height)

directory = "../results"
if not os.path.exists(directory):
    os.mkdir(directory)

filename = directory + os.path.sep + "AddImage.pptx"
prs.save(filename)
