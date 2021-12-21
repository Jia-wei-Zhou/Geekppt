from pptx import Presentation
import os


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello python-pptx"
subtitle.text = "python pptx test"

directory = "../results"
if not os.path.exists(directory):
    os.mkdir(directory)

filename = directory + os.path.sep + "HelloWorld.pptx"
prs.save(filename)
