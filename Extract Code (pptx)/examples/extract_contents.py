from pptx import Presentation
import os


input_path = "../resources" + os.path.sep + "test.pptx"
if not os.path.exists(input_path):
    print("{} does not exits".format(input_path))
    exit(1)

prs = Presentation(input_path)

texts = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                texts.append(run.text)

print(texts)
