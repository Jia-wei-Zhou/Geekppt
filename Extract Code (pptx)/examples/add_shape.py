from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
import os

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = "Add an AutoShape"

left = Inches(0.93)
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = "Step 1"

left = left + width - Inches(0.4)
width = Inches(2.0)

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = "Step {}".format(n)
    left = left + width - Inches(0.4)

directory = "../results"
if not os.path.exists(directory):
    os.mkdir(directory)

filename = directory + os.path.sep + "AddShape.pptx"
prs.save(filename)
