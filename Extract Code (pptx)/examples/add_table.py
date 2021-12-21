from pptx import Presentation
from pptx.util import Inches
import os


prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = "Adding a Table"

row = col = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(row, col, left, top, width, height).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

table.cell(0, 0).text = "Foo"
table.cell(0, 1).text = "Bar"

table.cell(1, 0).text = "Baz"
table.cell(1, 1).text = "Qux"

directory = "../results"
if not os.path.exists(directory):
    os.mkdir(directory)

filename = directory + os.path.sep + "AddTable.pptx"
prs.save(filename)