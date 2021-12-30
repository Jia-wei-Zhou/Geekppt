from pptx import Presentation


class PresentationBuilder(object):
    def __init__(self, filename):
        self.filename = filename
        self.presentation = Presentation(filename)

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst  # pylint: disable=protected-access

    def move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    # also works for deleting slides
    def delete_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index])

    def add_slide(self, layout_num=1):
        self.presentation.slides.add_slide(self.presentation.slide_layouts[layout_num])

    def save(self, filename=''):
        self.presentation.save(filename if filename != '' else self.filename)
