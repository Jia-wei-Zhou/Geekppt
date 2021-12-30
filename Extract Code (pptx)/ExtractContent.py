from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pypandoc as pandoc
import os


class ExtractContent:
    def __init__(self, snippets, filename):
        if not os.path.exists(filename):
            print("{} does not exist".format(filename))
            exit(1)

        self.filename = filename
        self.snippets = snippets
        self.pre = Presentation(filename)

    def __convert_code(self, slide, shape, contents):
        # determine if there is code snippet in this shape
        found = False
        for content in contents:
            if content in self.snippets:
                found = True
        if not found:
            return

        # delete original shape
        sp = shape._sp
        sp.getparent().remove(sp)

        # create new text boxes for the original contents
        for i in range(len(contents)):
            left = top = width = height = Inches(1 + i * 2)
            if contents[i] in self.snippets:
                self.__add_textbox(slide, contents[i], left, top, width, height, 255, 0, 255)
            else:
                self.__add_textbox(slide, contents[i], left, top, width, height, 0, 255, 255)

    @staticmethod
    def __add_textbox(slide, content, left, top, width, height, red, green, blue):
        # create new text box at specific position
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = content
        # set font for the content
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(red, green, blue)

    def find_match(self):
        # iterate through slides in the file
        for slide in self.pre.slides:
            # iterate through shapes in the slide
            for shape in slide.shapes:
                contents = []
                if not shape.has_text_frame:
                    continue
                # iterate through paragraphs in the shape
                for paragraph in shape.text_frame.paragraphs:
                    contents.append(paragraph.text)
            self.__convert_code(slide, shape, contents)

    def save(self, filename=''):
        self.pre.save(filename if filename != '' else self.filename)


# find code snippet from md file
def find_snippet(filename):
    if not os.path.exists(filename):
        print("{} does not exist".format(filename))
        exit(1)

    snippet = []
    with open(filename) as file:
        line = file.readline()
        while line:
            if line == "```\n":
                temp = ""
                line = file.readline()
                while line != "```\n":
                    temp += line
                    line = file.readline()
                snippet.append(temp)
            line = file.readline()

    return [piece.strip() for piece in snippet]


# convert md file to ppt slides 
def convert_file(filename):
    if not os.path.exists(filename):
        print("{} does not exist".format(filename))
        exit(1)

    out_file = filename.split('.')[0] + ".pptx"
    pandoc.convert_file(filename, "pptx", outputfile=out_file)


if __name__ == '__main__':
    filename = "resources/test.md"
    snippets = find_snippet(filename)
    convert_file(filename)

    extract = ExtractContent(snippets, "resources/test.pptx")
    extract.find_match()
    extract.save()
